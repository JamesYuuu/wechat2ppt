import collections.abc
import pptx
import os
from PIL import Image
import requests
from io import BytesIO
import shutil


class ppt(object):
    def __init__(self,link:list,width:str,ratio:str):
        self.link=link
        self.width=width
        self.ratio=ratio

        print(self.width,self.ratio)
    
    def save_image(self,img_src,id):
        image=Image.open(BytesIO(requests.get(img_src).content))
        image=image.resize((int(1080/float(self.ratio)),1080),resample=Image.Resampling.BICUBIC)
        file='./image/'+str(id)+'.png'
        image.save(file)

    def download_image(self):
        for data in self.link:
            if data['width']==self.width and data['ratio']==self.ratio:
                self.save_image(data['url'],data['id'])

    def make_ppt(self):

        ppt=pptx.Presentation()
    
        # 默认比例为4:3，根据图片大小把修改ppt尺寸  
        ppt.slide_width=int(ppt.slide_height/float(self.ratio))

        if os.path.exists('./image'):
            shutil.rmtree('./image')
        os.mkdir('./image')

        self.download_image()

        for file in sorted(os.listdir('./image'),key=lambda x:int(x.split('.')[0])):
            slide = ppt.slides.add_slide(ppt.slide_layouts[3])
            slide.shapes.add_picture(f'./image/{file}',0,0,ppt.slide_width,ppt.slide_height)

        ppt.save('./image/result.pptx')